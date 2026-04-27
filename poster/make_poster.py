"""HAAG conference-poster generator — "Field Guide".

24 x 36 portrait, all native python-pptx shapes. Design inspired by
natural-history field guides and the McGuire 2025 GRC poster:
asymmetric title with dominant hero figure, six numbered sections with
parchment ribbon tags + INDIGO circled numerals hanging off the left,
a body-mass strata strip running full width, a results row wired with
large TERRACOTTA '+' and '=' symbols (visual equation), and silhouette
placeholders in otherwise empty margin areas.

Run from the repo root:

    python poster/make_poster.py
"""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from palette import (
    CLAY_RGB, FOG_RGB, FOREST_RGB, INDIGO_RGB, IRON_RGB, PALETTE_ROLES,
    PAPER_RGB, PARCHMENT_RGB, SKY_RGB, STONE_RGB, STRAW_RGB,
    TERRACOTTA_RGB, WHITE_RGB,
)

__version__ = "0.4.0"


# ============================================================
# Conference metadata — EDIT THIS BLOCK for a new conference.
# ============================================================

OUTPUT_FILENAME = "poster.pptx"

CONFERENCE_TAG = "OMSCS Conference 2026"

AFFILIATION = ("Georgia Institute of Technology \u00b7 "
               "Human-Augmented Analytics Group")

FOOTER_LEFT   = "Human-Augmented Analytics Group \u00b7 Georgia Tech"
FOOTER_CENTER = ("OMSCS Conference 2026 \u00b7 Poster & Demo Session "
                 "\u00b7 May 12, 2026")
FOOTER_RIGHT  = "[[FILL: presenter email]]"


# ============================================================
# Page + layout
# ============================================================

PAGE_W, PAGE_H = 24.0, 36.0
MARGIN = 0.60

# Fonts
FONT_TITLE      = 66    # huge serif question title
FONT_SUBTITLE   = 20    # italic paper-title subtitle
FONT_AUTHOR     = 18
FONT_ADVISOR    = 14
FONT_AFFIL      = 13
FONT_SECTION    = 26    # ribbon-tag titles
FONT_NUMERAL    = 44
FONT_BODY       = 17
FONT_KEYFIND    = 32
FONT_PULLQ      = 22
FONT_CAPTION    = 13
FONT_FOOTER     = 12
FONT_STRATA     = 18    # strata strip labels
FONT_EQUATION   = 72    # + / = signs in results row
FONT_PILL       = 13

SERIF = "Georgia"
SANS  = "Arial"

BUILD_STAMP_COLOR = (180, 180, 180)


# ============================================================
# Registries
# ============================================================

PLACEHOLDERS: list[tuple[str, float, float]] = []


def _register(label: str, x: float, y: float) -> None:
    PLACEHOLDERS.append((label, round(x, 3), round(y, 3)))


# ============================================================
# Low-level helpers
# ============================================================

def _rgb(t):
    return RGBColor(*t)


def _set_no_line(shape):
    shape.line.fill.background()


def _set_line(shape, color_rgb, width_pt, dashed=False):
    line = shape.line
    line.color.rgb = _rgb(color_rgb)
    line.width = Pt(width_pt)
    if dashed:
        ln = line._get_or_add_ln()
        for child in list(ln):
            if child.tag == qn("a:prstDash"):
                ln.remove(child)
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")


def add_rect(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None,
             line_w_pt=1.0, dashed=False,
             shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_rgb is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_rgb)
    if line_rgb is None:
        _set_no_line(shape)
    else:
        _set_line(shape, line_rgb, line_w_pt, dashed=dashed)
    return shape


def add_oval(slide, x, y, w, h, *, fill_rgb=None, line_rgb=None,
             line_w_pt=1.0):
    return add_rect(slide, x, y, w, h, fill_rgb=fill_rgb,
                    line_rgb=line_rgb, line_w_pt=line_w_pt,
                    shape_type=MSO_SHAPE.OVAL)


def add_text(slide, x, y, w, h, text, *, face=SANS, size_pt=FONT_BODY,
             bold=False, italic=False, color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=True):
    box = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.04)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = face
    f.size = Pt(size_pt)
    f.bold = bold
    f.italic = italic
    f.color.rgb = _rgb(color_rgb)
    if register and text.startswith("[[") and text.endswith("]]"):
        _register(text, x, y)
    return box


def add_rule(slide, x, y, w, color_rgb, thickness_in=0.03):
    return add_rect(slide, x, y, w, thickness_in,
                    fill_rgb=color_rgb, line_rgb=None)


def add_placeholder(slide, x, y, w, h, label, *, size_pt=20,
                    label_color=STONE_RGB, border_color=IRON_RGB,
                    border_w_pt=1.0, fill=PAPER_RGB, dashed=True):
    add_rect(slide, x, y, w, h,
             fill_rgb=fill, line_rgb=border_color,
             line_w_pt=border_w_pt, dashed=dashed)
    add_text(slide, x, y, w, h, label,
             face=SANS, size_pt=size_pt, italic=True,
             color_rgb=label_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             register=False)
    _register(label, x, y)


def add_silhouette(slide, x, y, w, h, label="[[SILHOUETTE]]"):
    """Margin decoration — dashed FOG rect with italic label."""
    add_rect(slide, x, y, w, h,
             fill_rgb=None, line_rgb=FOG_RGB,
             line_w_pt=0.75, dashed=True)
    add_text(slide, x, y, w, h, label,
             face=SANS, size_pt=9, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             register=False)


# ============================================================
# Section "ribbon tag" — numbered badge + parchment title bar
# ============================================================

NUMERAL_CHARS = {1: "\u2460", 2: "\u2461", 3: "\u2462",
                 4: "\u2463", 5: "\u2464", 6: "\u2465"}
NUMERAL_DIAM   = 1.10
TAG_H          = 0.95
TAG_OVERLAP    = 0.55   # numeral overlaps into the tag by this much


def add_ribbon_tag(slide, x, y, w, n, title) -> float:
    """McGuire-style hanging tag: INDIGO circled numeral on the left
    overlapping a PARCHMENT ribbon carrying the section title.

    Returns the y-coordinate of the next free vertical position
    (bottom of the tag + small gap).
    """
    # Parchment ribbon — starts a bit right of the numeral's center so
    # the numeral overlaps it on the left.
    tag_x = x + NUMERAL_DIAM - TAG_OVERLAP
    tag_w = w - (tag_x - x)
    add_rect(slide, tag_x, y + (NUMERAL_DIAM - TAG_H) / 2,
             tag_w, TAG_H,
             fill_rgb=PARCHMENT_RGB, line_rgb=IRON_RGB, line_w_pt=0.75)

    # Title text inside the parchment
    add_text(slide, tag_x + 0.90, y + (NUMERAL_DIAM - TAG_H) / 2,
             tag_w - 1.00, TAG_H, title,
             face=SERIF, size_pt=FONT_SECTION, bold=True,
             color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # INDIGO filled circle + white numeral, drawn last so it sits on top
    add_oval(slide, x, y, NUMERAL_DIAM, NUMERAL_DIAM,
             fill_rgb=INDIGO_RGB, line_rgb=IRON_RGB, line_w_pt=1.25)
    add_text(slide, x, y, NUMERAL_DIAM, NUMERAL_DIAM, NUMERAL_CHARS[n],
             face=SERIF, size_pt=FONT_NUMERAL, bold=True,
             color_rgb=WHITE_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             register=False)

    return y + NUMERAL_DIAM + 0.15


# ============================================================
# Body helpers
# ============================================================

def add_body(slide, x, y, w, text, *, h=1.80,
             size_pt=FONT_BODY) -> float:
    add_text(slide, x, y, w, h, text,
             face=SANS, size_pt=size_pt, color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    return y + h


def add_bullets(slide, x, y, w, bullets, *, line_h=0.55,
                size_pt=FONT_BODY, bullet_rgb=TERRACOTTA_RGB) -> float:
    marker_w = 0.12
    for i, b in enumerate(bullets):
        by = y + i * line_h
        # small IRON-outlined TERRACOTTA square marker
        add_rect(slide, x, by + 0.14, marker_w, marker_w,
                 fill_rgb=bullet_rgb, line_rgb=None)
        add_text(slide, x + marker_w + 0.10, by,
                 w - marker_w - 0.12, line_h, b,
                 face=SANS, size_pt=size_pt, color_rgb=IRON_RGB,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 register=False)
        if b.startswith("[[") and b.endswith("]]"):
            _register(b, x, by)
    return y + len(bullets) * line_h


# ============================================================
# Page composers
# ============================================================

def build_background(slide):
    add_rect(slide, 0, 0, PAGE_W, PAGE_H, fill_rgb=PAPER_RGB)


def build_header(slide):
    """Asymmetric header — left: affiliation + authors; right: BIG title."""
    # Top thin rule
    add_rule(slide, MARGIN, 0.30, PAGE_W - 2 * MARGIN, IRON_RGB,
             thickness_in=0.04)

    # LEFT column: institution + authors + advisors + pill
    left_x = MARGIN
    add_text(slide, left_x, 0.55, 9.50, 0.35,
             "GEORGIA TECH \u00b7 HUMAN-AUGMENTED ANALYTICS GROUP",
             face=SERIF, size_pt=FONT_AFFIL, bold=True,
             color_rgb=INDIGO_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    add_text(slide, left_x, 0.95, 10.50, 0.45,
             "Spatial Camera-Trap Audit",
             face=SERIF, size_pt=22, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    add_text(slide, left_x, 1.55, 10.50, 0.35,
             "Cliff [[FILL: surname]]  \u00b7  Kefei Yan  \u00b7  "
             "Jacquie Carroll",
             face=SANS, size_pt=FONT_AUTHOR,
             color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_text(slide, left_x, 1.90, 10.50, 0.35,
             "Neelima Pandey  \u00b7  John Hiedo",
             face=SANS, size_pt=FONT_AUTHOR,
             color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    add_text(slide, left_x, 2.35, 10.50, 0.32,
             "Advisors:  Jenny McGuire  \u00b7  Steve Mussmann  "
             "\u00b7  Emily Weigel",
             face=SANS, size_pt=FONT_ADVISOR, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Conference pill — bottom of left column
    pill_w, pill_h = 3.40, 0.55
    pill_x, pill_y = left_x, 2.80
    add_rect(slide, pill_x, pill_y, pill_w, pill_h,
             fill_rgb=TERRACOTTA_RGB, line_rgb=None)
    add_text(slide, pill_x, pill_y, pill_w, pill_h, CONFERENCE_TAG,
             face=SERIF, size_pt=FONT_PILL, bold=True,
             color_rgb=WHITE_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # RIGHT column: huge question title
    title_x = 11.80
    title_w = PAGE_W - title_x - MARGIN
    add_text(slide, title_x, 0.50, title_w, 2.20,
             "Do cameras and range maps see the same biodiversity?",
             face=SERIF, size_pt=FONT_TITLE, bold=True,
             color_rgb=IRON_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle — italic, INDIGO
    add_text(slide, title_x, 2.80, title_w, 0.55,
             "Bridging Empirical and Modeled Biodiversity:  "
             "Snapshot USA \u00d7 IUCN \u00d7 COMBINE",
             face=SERIF, size_pt=FONT_SUBTITLE, italic=True,
             color_rgb=INDIGO_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Bottom thick rule
    add_rule(slide, MARGIN, 3.60, PAGE_W - 2 * MARGIN, IRON_RGB,
             thickness_in=0.05)


def build_top_row(slide):
    """Top band: ① Motivation (left) | Hero figure (right)."""
    y_top = 3.95

    # LEFT: (1) Motivation
    col_w = 11.20
    x = MARGIN
    y = add_ribbon_tag(slide, x, y_top, col_w, 1, "Motivation")
    y = add_body(
        slide, x + 0.30, y + 0.05, col_w - 0.60,
        "Across the U.S., camera traps record mammals through "
        "Snapshot USA (SSUSA). IUCN Red List maps \u2014 drawn "
        "largely from expert opinion \u2014 are the global standard "
        "for species ranges. The two views of biodiversity are rarely "
        "checked against each other, so misalignments persist "
        "unnoticed in conservation planning.",
        h=3.00, size_pt=FONT_BODY)

    # Pull-quote — the research question, big italic serif on the left
    add_rect(slide, x + 0.30, y + 0.15, 0.08, 1.70,
             fill_rgb=TERRACOTTA_RGB, line_rgb=None)
    add_text(slide, x + 0.60, y + 0.05, col_w - 0.90, 1.90,
             "\u201cWhen maps and cameras disagree, which one "
             "should we trust?\u201d",
             face=SERIF, size_pt=FONT_PULLQ, italic=True,
             color_rgb=INDIGO_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # RIGHT: Hero figure
    hero_x = MARGIN + col_w + 0.50
    hero_w = PAGE_W - hero_x - MARGIN
    add_placeholder(slide, hero_x, y_top, hero_w, 8.10,
                    "[[HERO VISUAL \u2014 SSUSA \u2229 IUCN across "
                    "the contiguous U.S. + AK + HI]]",
                    size_pt=18,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.5,
                    fill=PAPER_RGB,
                    dashed=False)
    add_text(slide, hero_x, y_top + 8.15, hero_w, 0.35,
             "Hero: cross-source agreement mapped by ecoregion.",
             face=SERIF, size_pt=FONT_CAPTION, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=False)


def build_data_sources_band(slide):
    """② Data Sources — compact strip below motivation on the left."""
    y_top = 9.60
    col_w = 11.20
    x = MARGIN
    y = add_ribbon_tag(slide, x, y_top, col_w, 2, "Data Sources")

    # 3 columns of data-source summaries, each with GOLDEN-highlighted name
    sub_w = (col_w - 0.60 - 0.50) / 3
    sub_y = y + 0.10
    sources = [
        ("SSUSA", "713,319 records",
         "109 species  \u00b7  2019\u20132024"),
        ("IUCN",  "733 polygons",
         "575 species  \u00b7  Red List"),
        ("COMBINE", "6,110 species",
         "Body-mass traits"),
    ]
    for i, (name, primary, secondary) in enumerate(sources):
        sx = x + 0.30 + i * (sub_w + 0.25)
        # Tight underline under the source name in INDIGO
        add_text(slide, sx, sub_y, sub_w, 0.35, name,
                 face=SERIF, size_pt=20, bold=True,
                 color_rgb=INDIGO_RGB,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 register=False)
        add_rule(slide, sx, sub_y + 0.38, 0.90, TERRACOTTA_RGB,
                 thickness_in=0.03)
        add_text(slide, sx, sub_y + 0.48, sub_w, 0.35, primary,
                 face=SANS, size_pt=17, bold=True,
                 color_rgb=IRON_RGB,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 register=False)
        add_text(slide, sx, sub_y + 0.85, sub_w, 0.30, secondary,
                 face=SANS, size_pt=12,
                 color_rgb=STONE_RGB,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 register=False)


def build_strata_strip(slide):
    """③ Body-Mass Strata — full-width timeline-style strip."""
    y_top = 13.00
    y = add_ribbon_tag(slide, MARGIN, y_top,
                       PAGE_W - 2 * MARGIN, 3,
                       "Body-Mass Strata \u2014 where SAC is computed")

    # 4 boxes in a horizontal strip with arrows between
    strip_y = y + 0.15
    strip_h = 2.00
    strips = [
        ("\u2264 50 g",     "small-bodied",      SKY_RGB),
        ("50\u2013100 g",    "small-medium",      FOREST_RGB),
        ("100\u2013500 g",   "medium",            STRAW_RGB),
        ("500\u20131000 g",  "medium-large",      CLAY_RGB),
        ("\u2265 1000 g",    "large-bodied",      TERRACOTTA_RGB),
    ]
    n = len(strips)
    arrow_w = 0.50
    avail_w = PAGE_W - 2 * MARGIN - (n - 1) * arrow_w
    box_w = avail_w / n
    for i, (label, sub, color) in enumerate(strips):
        bx = MARGIN + i * (box_w + arrow_w)
        # Colored tab
        add_rect(slide, bx, strip_y, box_w, 0.40,
                 fill_rgb=color, line_rgb=None)
        # White bordered body
        add_rect(slide, bx, strip_y + 0.40, box_w, strip_h - 0.40,
                 fill_rgb=PAPER_RGB, line_rgb=IRON_RGB, line_w_pt=1.0)
        # Threshold label — bold serif, centered
        add_text(slide, bx, strip_y + 0.50, box_w, 0.60, label,
                 face=SERIF, size_pt=22, bold=True,
                 color_rgb=IRON_RGB,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 register=False)
        # Sub-label — italic, STONE
        add_text(slide, bx, strip_y + 1.05, box_w, 0.40, sub,
                 face=SANS, size_pt=13, italic=True,
                 color_rgb=STONE_RGB,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                 register=False)
        # Silhouette placeholder at the bottom of each box
        sil_w = box_w * 0.6
        sil_x = bx + (box_w - sil_w) / 2
        add_silhouette(slide, sil_x, strip_y + 1.45, sil_w, 0.45,
                       "[[SILHOUETTE]]")
        # Arrow between boxes
        if i < n - 1:
            ax = bx + box_w + 0.05
            add_text(slide, ax, strip_y + 0.40, arrow_w - 0.10,
                     strip_h - 0.40, "\u2192",
                     face=SERIF, size_pt=32, bold=True,
                     color_rgb=INDIGO_RGB,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                     register=False)


def build_results_row(slide):
    """④ Results — visual equation: [fig] + [fig] = [headline]."""
    y_top = 16.20
    y = add_ribbon_tag(slide, MARGIN, y_top,
                       PAGE_W - 2 * MARGIN, 4, "Results")

    # Equation layout: 3 boxes separated by + and = symbols
    # Box A + Box B = Box C (headline)
    y_row = y + 0.15
    row_h = 5.80
    usable_w = PAGE_W - 2 * MARGIN
    sym_w = 0.90
    # Box A and B are equal width; Box C is larger (the result)
    small_box_w = (usable_w - 2 * sym_w - 0.60) / 3.2
    big_box_w = small_box_w * 1.4

    x_a = MARGIN
    x_plus = x_a + small_box_w + 0.10
    x_b = x_plus + sym_w + 0.10
    x_eq = x_b + small_box_w + 0.10
    x_c = x_eq + sym_w + 0.10

    # Box A: SAC completeness
    add_placeholder(slide, x_a, y_row, small_box_w, row_h - 0.90,
                    "[[FIGURE \u2014 SAC completeness curves by strata]]",
                    size_pt=14,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.25,
                    fill=PAPER_RGB,
                    dashed=False)
    add_text(slide, x_a, y_row + row_h - 0.85, small_box_w, 0.80,
             "Species-accumulation curves \u2014 when does an array "
             "look complete?",
             face=SERIF, size_pt=FONT_CAPTION, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=False)

    # + symbol
    add_text(slide, x_plus, y_row, sym_w, row_h - 0.90, "+",
             face=SERIF, size_pt=FONT_EQUATION, bold=True,
             color_rgb=TERRACOTTA_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             register=False)

    # Box B: CI-gap classifier
    add_placeholder(slide, x_b, y_row, small_box_w, row_h - 0.90,
                    "[[FIGURE \u2014 CI-gap upper bound by effort]]",
                    size_pt=14,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.25,
                    fill=PAPER_RGB,
                    dashed=False)
    add_text(slide, x_b, y_row + row_h - 0.85, small_box_w, 0.80,
             "Binomial CI-gap \u2014 how many species could still be missing?",
             face=SERIF, size_pt=FONT_CAPTION, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=False)

    # = symbol
    add_text(slide, x_eq, y_row, sym_w, row_h - 0.90, "=",
             face=SERIF, size_pt=FONT_EQUATION, bold=True,
             color_rgb=TERRACOTTA_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             register=False)

    # Box C: Headline — TWILIGHT background with big finding text on top,
    # figure below.
    head_w = usable_w - (x_c - MARGIN)
    kf_h = 2.00
    add_rect(slide, x_c, y_row, head_w, kf_h,
             fill_rgb=INDIGO_RGB, line_rgb=None)
    add_text(slide, x_c + 0.30, y_row, head_w - 0.60, kf_h,
             "SAC alone misclassifies 4.4% "
             "of array-years as adequate.",
             face=SERIF, size_pt=FONT_KEYFIND, bold=True,
             color_rgb=PAPER_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             register=False)
    # Figure under the finding
    add_placeholder(slide, x_c, y_row + kf_h + 0.15,
                    head_w, row_h - kf_h - 1.05,
                    "[[HEADLINE FIGURE \u2014 2\u00d72 SAC \u00d7 CI-gap]]",
                    size_pt=16,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.25,
                    fill=PAPER_RGB,
                    dashed=False)
    add_text(slide, x_c, y_row + row_h - 0.85, head_w, 0.80,
             "False-plateau array-years \u2014 apparent "
             "completeness but missed species.",
             face=SERIF, size_pt=FONT_CAPTION, italic=True,
             color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=False)


def build_bottom_row(slide):
    """⑤ Supporting (left) | ⑥ Takeaways (right)."""
    y_top = 23.00
    col_w = (PAGE_W - 2 * MARGIN - 0.60) / 2

    # LEFT: (5) Supporting Analysis
    x = MARGIN
    y = add_ribbon_tag(slide, x, y_top, col_w, 5,
                       "Supporting Analysis")
    y = add_body(
        slide, x + 0.30, y + 0.05, col_w - 0.60,
        "Community composition is most stable for mammals above "
        "100 g. Below 50 g, sparse detections drive high variance "
        "in both accumulation curves and IUCN overlap. Multi-year "
        "rollups (2019\u20132024) dampen single-year noise. "
        "Calendar-night effort \u2014 not camera-night, which "
        "inflates ~10\u00d7 \u2014 is the correct denominator when "
        "pooling across arrays.",
        h=3.40, size_pt=FONT_BODY)
    y += 0.15
    add_placeholder(slide, x + 0.30, y, col_w - 0.60, 3.40,
                    "[[SUPPORTING FIGURE]]",
                    size_pt=16,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.25,
                    fill=PAPER_RGB,
                    dashed=False)

    # RIGHT: (6) Takeaways & Next Steps
    x = MARGIN + col_w + 0.60
    y = add_ribbon_tag(slide, x, y_top, col_w, 6,
                       "Takeaways & Next Steps")
    y = add_bullets(slide, x + 0.30, y + 0.05, col_w - 0.60, [
        "Empirical and modeled biodiversity views mostly agree "
        "for mid-sized mammals.",
        "SAC alone under-reports sampling gaps; the CI-gap "
        "classifier catches false plateaus.",
        "Cross-source validation is a cheap, reproducible audit "
        "for any camera-trap benchmark.",
        "Next: extend to body-mass and guild-level stratified "
        "audits across all SSUSA years.",
    ], line_h=1.10, size_pt=FONT_BODY)
    y += 0.30

    # Repo link (bold INDIGO)
    add_text(slide, x + 0.30, y, col_w - 1.80, 0.35,
             "github.com / cliff003/ HAAG_Spatial_Camera_Trap_Fall2025",
             face=SANS, size_pt=14, bold=True,
             color_rgb=INDIGO_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             register=False)

    # QR placeholder — bottom-right of this panel
    qr_size = 1.30
    qr_x = x + col_w - qr_size - 0.30
    qr_y = y_top + 6.30
    add_placeholder(slide, qr_x, qr_y, qr_size, qr_size,
                    "[[QR]]", size_pt=11,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.0,
                    fill=PAPER_RGB,
                    dashed=True)


def build_silhouettes_and_illustration(slide):
    """Scatter silhouettes in empty margin areas + a field-plate strip."""
    # Left margin silhouette near motivation
    add_silhouette(slide, 0.15, 5.40, 0.40, 0.55,
                   "[[SILH]]")
    # Left margin silhouette near data sources
    add_silhouette(slide, 0.15, 10.80, 0.40, 0.55,
                   "[[SILH]]")
    # Right margin silhouette beside headline
    add_silhouette(slide, 23.45, 18.50, 0.40, 0.55,
                   "[[SILH]]")
    # Bottom-right vintage illustration — larger, field-guide plate
    add_placeholder(slide, 18.60, 30.50, 4.80, 2.60,
                    "[[FIELD-PLATE ILLUSTRATION]]",
                    size_pt=13,
                    label_color=STONE_RGB,
                    border_color=IRON_RGB,
                    border_w_pt=1.0,
                    fill=PARCHMENT_RGB,
                    dashed=True)


def build_footer(slide):
    """Thin footer with IRON rule, swatch strip, meta, build stamp."""
    y = 33.50
    add_rule(slide, MARGIN, y, PAGE_W - 2 * MARGIN, IRON_RGB,
             thickness_in=0.04)
    y += 0.12

    # Swatch strip — centered
    add_swatch_strip(slide, y_row=y)
    y += 0.80

    # Meta text
    y_txt = y + 0.10
    add_text(slide, MARGIN, y_txt, 8.00, 0.30, FOOTER_LEFT,
             face=SANS, size_pt=FONT_FOOTER, color_rgb=STONE_RGB,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    add_text(slide, 8.00, y_txt, 8.00, 0.30, FOOTER_CENTER,
             face=SANS, size_pt=FONT_FOOTER, color_rgb=STONE_RGB,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_text(slide, 16.00, y_txt, PAGE_W - 16.00 - MARGIN, 0.30,
             FOOTER_RIGHT,
             face=SANS, size_pt=FONT_FOOTER, color_rgb=STONE_RGB,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)

    # Build stamp
    stamp = f"v{__version__} \u00b7 {datetime.now():%Y-%m-%d %H:%M}"
    add_text(slide, PAGE_W - 4.00, PAGE_H - 0.35, 3.50, 0.25, stamp,
             face=SANS, size_pt=8, color_rgb=BUILD_STAMP_COLOR,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM,
             register=False)


def add_swatch_strip(slide, y_row):
    """11 swatches centered, with hex labels below."""
    n = len(PALETTE_ROLES)
    sw_w, sw_h, gap = 0.55, 0.30, 0.05
    total_w = n * sw_w + (n - 1) * gap
    x_start = (PAGE_W - total_w) / 2
    for i, (_, hex_, rgb_t, _) in enumerate(PALETTE_ROLES):
        x = x_start + i * (sw_w + gap)
        outline = IRON_RGB if rgb_t == PAPER_RGB else None
        add_rect(slide, x, y_row, sw_w, sw_h,
                 fill_rgb=rgb_t, line_rgb=outline, line_w_pt=0.5)
        add_text(slide, x - 0.10, y_row + sw_h + 0.02, sw_w + 0.20,
                 0.15, hex_,
                 face=SANS, size_pt=7, color_rgb=STONE_RGB,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                 register=False)


# ============================================================
# Main + summary
# ============================================================

def print_build_summary(out_path: Path):
    stamp = f"v{__version__} \u00b7 {datetime.now():%Y-%m-%d %H:%M}"
    size_kb = out_path.stat().st_size / 1024
    print("=" * 72)
    print(f"HAAG poster build summary  \u2014  {CONFERENCE_TAG}")
    print("Design: \"Field Guide\" (McGuire-inspired)  \u00b7  "
          "palette: \"Field Guide\"")
    print("=" * 72)
    print(f"Output:      {out_path}")
    print(f"Slide dims:  {PAGE_W} x {PAGE_H} in (portrait)")
    print(f"Build stamp: {stamp}")
    print(f"File size:   {size_kb:.1f} KB")
    print()
    print("Palette (11 tokens):")
    for name, hex_, rgb_t, role in PALETTE_ROLES:
        print(f"  {name:10} {hex_}  rgb{str(rgb_t):>17}   {role}")
    print()
    print("Layout:")
    print("  Header      y  0.00 \u2013  3.65  (asymmetric: meta | big title)")
    print("  \u2460 Motivation + Hero  y  3.95 \u2013  9.50")
    print("  \u2461 Data Sources       y  9.60 \u2013 12.70")
    print("  \u2462 Body-Mass Strata   y 13.00 \u2013 16.00 (full width)")
    print("  \u2463 Results (eqn)      y 16.20 \u2013 22.70 (full width)")
    print("  \u2464 Supporting         y 23.00 \u2013 30.50 (left 1/2)")
    print("  \u2465 Takeaways          y 23.00 \u2013 30.50 (right 1/2)")
    print("  Footer                y 33.50 \u2013 36.00")
    print()
    print(f"Placeholders: {len(PLACEHOLDERS)}")
    for label, x, y in PLACEHOLDERS:
        print(f"  ({x:5.2f}, {y:5.2f})  {label}")
    print("=" * 72)


def main() -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    assert prs.slide_width  == Inches(PAGE_W)
    assert prs.slide_height == Inches(PAGE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    build_background(slide)
    build_header(slide)
    build_top_row(slide)
    build_data_sources_band(slide)
    build_strata_strip(slide)
    build_results_row(slide)
    build_bottom_row(slide)
    build_silhouettes_and_illustration(slide)
    build_footer(slide)

    out = Path(__file__).parent / OUTPUT_FILENAME
    prs.save(out)

    size = out.stat().st_size
    assert size < 300_000, f"pptx unexpectedly large: {size} bytes"

    print_build_summary(out)
    return out


if __name__ == "__main__":
    main()
