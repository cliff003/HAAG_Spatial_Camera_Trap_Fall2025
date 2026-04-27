"""Drop notebook-generated figures into poster.pptx placeholders.

For each FIGURE_MAP entry, removes the dashed `[[...]]` placeholder from a
fresh copy of `poster/poster.pptx` and inserts the matching PNG at the
same (x, y, w, h). Saves to a named output PPTX.

Run from the repo root:

    # Default: build 3 preview posters (one per hero variant)
    python poster/merge_figures.py --batch

    # Single custom build
    python poster/merge_figures.py --hero H1 --boxB B3 --out poster/custom.pptx
"""
from __future__ import annotations

import argparse
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

SRC_PPTX = Path("poster/poster.pptx")
PNG_DIR = Path("figures/poster/png")


# Placeholders present in the empty poster, mapped to *label fragment* used
# to locate the shape in poster.pptx (from make_poster.py output).
SLOT_TO_LABEL = {
    "hero":        "HERO VISUAL",
    "boxA":        "FIGURE \u2014 SAC",
    "boxB":        "FIGURE \u2014 CI-gap",   # retained label in template
    "boxC":        "HEADLINE FIGURE",
    "supporting":  "SUPPORTING FIGURE",
    "fieldplate":  "FIELD-PLATE ILLUSTRATION",
}

# Default assignments per variant choice. Any slot missing here is skipped.
DEFAULT_BUILDS = {
    "H1": {
        "hero":        "H1_hero_out_of_range.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H2": {
        "hero":        "H2_hero_buffer_coverage.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H3": {
        "hero":        "H3_hero_combined.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H4a": {
        "hero":        "H4a_hero_elk_rawdata.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H4b": {
        "hero":        "H4b_hero_moose_rawdata.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H4c": {
        "hero":        "H4c_hero_blackbear_rawdata.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H4d": {
        "hero":        "H4d_hero_wtdeer_rawdata.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
    "H4e": {
        "hero":        "H4e_hero_raccoon_rawdata.png",
        "boxA":        "V6_boxA_sac_by_mass.png",
        "boxB":        "B3_boxB_jaccard_effort_dual.png",
        "boxC":        "C1_boxC_jaccard_effort_by_habitat.png",
        "supporting":  "S1_supporting_rf_importance.png",
        "fieldplate":  "V5_fieldplate_flagship_drilldown.png",
    },
}


def _iter_text_shapes(slide):
    for sh in list(slide.shapes):
        if sh.has_text_frame:
            yield sh, sh.text_frame.text


def _find_placeholder(slide, label_fragment):
    """Return (label_shape, rect_shape, (x, y, w, h)) for a placeholder."""
    label_shape = None
    for sh, txt in _iter_text_shapes(slide):
        if label_fragment in txt:
            label_shape = sh
            break
    if label_shape is None:
        return None
    x, y, w, h = (label_shape.left, label_shape.top,
                  label_shape.width, label_shape.height)
    rect_shape = None
    for sh in list(slide.shapes):
        if sh is label_shape:
            continue
        if (sh.left == x and sh.top == y
                and sh.width == w and sh.height == h):
            rect_shape = sh
            break
    return label_shape, rect_shape, (x, y, w, h)


def _delete_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def build(figure_map: dict, out_path: Path, verbose: bool = True):
    assert SRC_PPTX.exists(), f"missing {SRC_PPTX}"
    prs = Presentation(str(SRC_PPTX))
    slide = prs.slides[0]

    missing = [n for n in figure_map.values() if not (PNG_DIR / n).exists()]
    if missing:
        raise SystemExit(
            f"Missing PNGs in {PNG_DIR}: {missing}. "
            "Re-run Poster_Figures_OMSCS2026.ipynb first."
        )

    for slot, png_name in figure_map.items():
        label_frag = SLOT_TO_LABEL.get(slot)
        if not label_frag:
            print(f"  SKIP: unknown slot '{slot}'"); continue
        found = _find_placeholder(slide, label_frag)
        if found is None:
            print(f"  SKIP: placeholder '{label_frag}' not found"); continue
        label_shape, rect_shape, (x, y, w, h) = found
        _delete_shape(label_shape)
        if rect_shape is not None:
            _delete_shape(rect_shape)
        slide.shapes.add_picture(str(PNG_DIR / png_name), x, y, w, h)
        if verbose:
            print(f"  {slot:11s} <- {png_name}  at "
                  f"({x/914400:.2f},{y/914400:.2f}) "
                  f"{w/914400:.2f}x{h/914400:.2f}in")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"wrote {out_path}")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--batch", action="store_true",
                    help="Build one preview per hero variant (H1/H2/H3)")
    ap.add_argument("--hero",
                    choices=["H1", "H2", "H3",
                             "H4a", "H4b", "H4c", "H4d", "H4e"],
                    default="H1")
    ap.add_argument("--out", default="poster/poster_with_figures.pptx")
    args = ap.parse_args()

    if args.batch:
        for tag, fig_map in DEFAULT_BUILDS.items():
            out = Path(f"poster/poster_preview_{tag}.pptx")
            print(f"=== {tag} ===")
            build(fig_map, out)
    else:
        build(DEFAULT_BUILDS[args.hero], Path(args.out))


if __name__ == "__main__":
    main()
